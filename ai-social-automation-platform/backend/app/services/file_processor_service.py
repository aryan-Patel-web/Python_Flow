import os
import json
import logging
import mimetypes
from typing import Dict, List, Optional, Any, Union
from datetime import datetime
from pathlib import Path
import hashlib

# File processing imports
import PyPDF2
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd
from langdetect import detect
import textract

from config.database import db_instance
from utils.response_utils import create_response

logger = logging.getLogger(__name__)

class FileProcessorService:
    def __init__(self):
        self.supported_formats = {
            'pdf': ['.pdf'],
            'document': ['.docx', '.doc'],
            'presentation': ['.pptx', '.ppt'],
            'image': ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'],
            'text': ['.txt', '.md', '.rtf'],
            'spreadsheet': ['.xlsx', '.xls', '.csv'],
            'data': ['.json', '.xml']
        }
        
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.upload_folder = os.getenv('UPLOAD_FOLDER', 'uploads/temp')
        
        # Ensure upload directory exists
        os.makedirs(self.upload_folder, exist_ok=True)

    def process_file(self, file_id: str) -> Dict:
        """
        Process uploaded file and extract content for AI generation
        """
        try:
            db = db_instance.get_db()
            
            # Get file record
            file_record = db.uploaded_files.find_one({'_id': file_id})
            if not file_record:
                return {'success': False, 'error': 'File not found'}
            
            file_path = file_record['file_path']
            if not os.path.exists(file_path):
                return {'success': False, 'error': 'File not accessible'}
            
            # Detect file type
            file_type = self._detect_file_type(file_path)
            
            # Process based on file type
            if file_type == 'pdf':
                result = self._process_pdf(file_path)
            elif file_type == 'document':
                result = self._process_document(file_path)
            elif file_type == 'presentation':
                result = self._process_presentation(file_path)
            elif file_type == 'image':
                result = self._process_image(file_path)
            elif file_type == 'text':
                result = self._process_text_file(file_path)
            elif file_type == 'spreadsheet':
                result = self._process_spreadsheet(file_path)
            elif file_type == 'data':
                result = self._process_data_file(file_path)
            else:
                return {'success': False, 'error': 'Unsupported file type'}
            
            if result['success']:
                # Update file record with processed content
                db.uploaded_files.update_one(
                    {'_id': file_id},
                    {
                        '$set': {
                            'processed_at': datetime.utcnow(),
                            'content_extracted': True,
                            'processing_result': result
                        }
                    }
                )
            
            return result
            
        except Exception as e:
            logger.error(f"File processing error: {str(e)}")
            return {'success': False, 'error': f'Processing failed: {str(e)}'}

    def _detect_file_type(self, file_path: str) -> str:
        """Detect file type from extension"""
        extension = Path(file_path).suffix.lower()
        
        for file_type, extensions in self.supported_formats.items():
            if extension in extensions:
                return file_type
        
        return 'unknown'

    def _process_pdf(self, file_path: str) -> Dict:
        """Process PDF files"""
        try:
            content = ""
            metadata = {}
            
            # Try pdfplumber first (better for complex layouts)
            try:
                with pdfplumber.open(file_path) as pdf:
                    metadata = {
                        'pages': len(pdf.pages),
                        'title': pdf.metadata.get('Title', ''),
                        'author': pdf.metadata.get('Author', ''),
                        'subject': pdf.metadata.get('Subject', '')
                    }
                    
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n\n"
            except:
                # Fallback to PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata = {
                        'pages': len(pdf_reader.pages),
                        'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                        'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else ''
                    }
                    
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n\n"
            
            # Generate summary and context
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(content)
            
            return {
                'success': True,
                'content': content.strip(),
                'context': {
                    'file_type': 'text',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'word_count': len(content.split()),
                    'language': self._detect_language(content)
                }
            }
            
        except Exception as e:
            logger.error(f"Text file processing error: {str(e)}")
            return {'success': False, 'error': f'Text file processing failed: {str(e)}'}

    def _process_spreadsheet(self, file_path: str) -> Dict:
        """Process Excel/CSV files"""
        try:
            content_parts = []
            metadata = {}
            
            if file_path.endswith('.csv'):
                # Process CSV
                df = pd.read_csv(file_path)
                metadata = {
                    'format': 'csv',
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                }
                
                # Convert to readable format
                content_parts.append(f"Data Summary:\nRows: {len(df)}, Columns: {len(df.columns)}")
                content_parts.append(f"Column Names: {', '.join(df.columns.tolist())}")
                
                # Sample data (first 10 rows)
                if len(df) > 0:
                    content_parts.append("\nSample Data:")
                    content_parts.append(df.head(10).to_string(index=False))
                
                # Basic statistics for numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    content_parts.append("\nNumeric Column Statistics:")
                    content_parts.append(df[numeric_cols].describe().to_string())
                
            else:
                # Process Excel
                excel_data = pd.read_excel(file_path, sheet_name=None)
                sheet_names = list(excel_data.keys())
                
                metadata = {
                    'format': 'excel',
                    'sheets': len(sheet_names),
                    'sheet_names': sheet_names
                }
                
                content_parts.append(f"Excel Workbook Summary:\nSheets: {len(sheet_names)}")
                content_parts.append(f"Sheet Names: {', '.join(sheet_names)}")
                
                # Process each sheet
                for sheet_name, df in excel_data.items():
                    content_parts.append(f"\n--- Sheet: {sheet_name} ---")
                    content_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                    
                    if len(df.columns) > 0:
                        content_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                    
                    # Sample data (first 5 rows per sheet)
                    if len(df) > 0:
                        content_parts.append("Sample Data:")
                        content_parts.append(df.head(5).to_string(index=False))
            
            content = "\n".join(content_parts)
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(content)
            
            return {
                'success': True,
                'content': content,
                'context': {
                    'file_type': 'spreadsheet',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'data_structure': 'tabular',
                    'language': 'data'
                }
            }
            
        except Exception as e:
            logger.error(f"Spreadsheet processing error: {str(e)}")
            return {'success': False, 'error': f'Spreadsheet processing failed: {str(e)}'}

    def _process_data_file(self, file_path: str) -> Dict:
        """Process JSON/XML data files"""
        try:
            content = ""
            metadata = {}
            
            if file_path.endswith('.json'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    data = json.load(file)
                
                metadata = {
                    'format': 'json',
                    'structure': type(data).__name__,
                    'keys': list(data.keys()) if isinstance(data, dict) else None,
                    'length': len(data) if isinstance(data, (list, dict)) else None
                }
                
                # Convert to readable format
                content = json.dumps(data, indent=2, ensure_ascii=False)[:2000]  # Limit size
                if len(content) >= 2000:
                    content += "... (truncated)"
                
            elif file_path.endswith('.xml'):
                with open(file_path, 'r', encoding='utf-8') as file:
                    content = file.read()
                
                metadata = {
                    'format': 'xml',
                    'size': len(content)
                }
                
                # Truncate if too long
                if len(content) > 2000:
                    content = content[:2000] + "... (truncated)"
            
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(content)
            
            return {
                'success': True,
                'content': content,
                'context': {
                    'file_type': 'data',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'data_structure': 'structured',
                    'language': 'data'
                }
            }
            
        except Exception as e:
            logger.error(f"Data file processing error: {str(e)}")
            return {'success': False, 'error': f'Data file processing failed: {str(e)}'}

    def _generate_content_summary(self, content: str) -> str:
        """Generate a brief summary of content"""
        try:
            if not content.strip():
                return "No content available"
            
            # Simple summarization logic (can be enhanced with AI)
            sentences = content.replace('\n', ' ').split('.')
            sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
            
            if len(sentences) <= 3:
                return content[:200] + ("..." if len(content) > 200 else "")
            
            # Take first, middle, and last sentences
            summary_sentences = [
                sentences[0],
                sentences[len(sentences)//2] if len(sentences) > 2 else "",
                sentences[-1] if len(sentences) > 1 else ""
            ]
            
            summary = ". ".join([s for s in summary_sentences if s]).strip()
            
            # Limit length
            if len(summary) > 300:
                summary = summary[:300] + "..."
            
            return summary
            
        except Exception as e:
            logger.error(f"Summary generation error: {str(e)}")
            return "Summary generation failed"

    def _extract_key_points(self, content: str) -> List[str]:
        """Extract key points from content"""
        try:
            if not content.strip():
                return []
            
            # Simple key point extraction (can be enhanced with NLP)
            lines = content.split('\n')
            key_points = []
            
            # Look for bullet points, numbered lists, or headers
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Bullet points
                if line.startswith(('•', '-', '*', '▪', '▫')):
                    key_points.append(line[1:].strip())
                
                # Numbered lists
                elif line[0].isdigit() and '.' in line[:5]:
                    key_points.append(line.split('.', 1)[1].strip())
                
                # Headers (short lines that might be important)
                elif len(line) < 60 and not line.endswith('.') and ':' in line:
                    key_points.append(line)
                
                # Important phrases (containing keywords)
                elif any(keyword in line.lower() for keyword in [
                    'important', 'key', 'note', 'remember', 'conclusion',
                    'summary', 'main', 'primary', 'essential', 'critical'
                ]):
                    key_points.append(line[:100] + ("..." if len(line) > 100 else ""))
            
            # Limit to top 10 key points
            return key_points[:10]
            
        except Exception as e:
            logger.error(f"Key points extraction error: {str(e)}")
            return []

    def _detect_language(self, content: str) -> str:
        """Detect language of content"""
        try:
            if not content.strip() or len(content.strip()) < 10:
                return 'unknown'
            
            # Use first 500 characters for detection
            sample_text = content[:500]
            detected_lang = detect(sample_text)
            
            # Map to full language names
            lang_map = {
                'en': 'English',
                'es': 'Spanish',
                'fr': 'French',
                'de': 'German',
                'it': 'Italian',
                'pt': 'Portuguese',
                'ru': 'Russian',
                'zh': 'Chinese',
                'ja': 'Japanese',
                'ko': 'Korean',
                'ar': 'Arabic',
                'hi': 'Hindi'
            }
            
            return lang_map.get(detected_lang, detected_lang)
            
        except Exception as e:
            logger.debug(f"Language detection failed: {str(e)}")
            return 'unknown'

    def upload_file(self, user_id: str, file_data, filename: str) -> Dict:
        """Handle file upload and initial processing"""
        try:
            # Validate file
            if not file_data:
                return {'success': False, 'error': 'No file provided'}
            
            # Check file size
            file_data.seek(0, 2)  # Seek to end
            file_size = file_data.tell()
            file_data.seek(0)  # Reset to beginning
            
            if file_size > self.max_file_size:
                return {'success': False, 'error': f'File too large. Max size: {self.max_file_size // (1024*1024)}MB'}
            
            # Validate file type
            file_extension = Path(filename).suffix.lower()
            file_type = self._detect_file_type(filename)
            
            if file_type == 'unknown':
                return {'success': False, 'error': 'Unsupported file type'}
            
            # Generate unique filename
            timestamp = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
            file_hash = hashlib.md5(f"{user_id}_{timestamp}_{filename}".encode()).hexdigest()[:8]
            unique_filename = f"{timestamp}_{file_hash}_{filename}"
            
            # Save file
            file_path = os.path.join(self.upload_folder, unique_filename)
            file_data.save(file_path)
            
            # Create database record
            db = db_instance.get_db()
            file_record = {
                'user_id': user_id,
                'original_filename': filename,
                'unique_filename': unique_filename,
                'file_path': file_path,
                'file_size': file_size,
                'file_type': file_type,
                'mime_type': mimetypes.guess_type(filename)[0],
                'uploaded_at': datetime.utcnow(),
                'processed': False,
                'content_extracted': False
            }
            
            result = db.uploaded_files.insert_one(file_record)
            file_id = str(result.inserted_id)
            
            return {
                'success': True,
                'file_id': file_id,
                'filename': filename,
                'file_type': file_type,
                'file_size': file_size,
                'message': 'File uploaded successfully'
            }
            
        except Exception as e:
            logger.error(f"File upload error: {str(e)}")
            return {'success': False, 'error': f'Upload failed: {str(e)}'}

    def get_user_files(self, user_id: str, limit: int = 20, offset: int = 0) -> Dict:
        """Get user's uploaded files"""
        try:
            db = db_instance.get_db()
            
            files = list(
                db.uploaded_files
                .find({'user_id': user_id})
                .sort('uploaded_at', -1)
                .limit(limit)
                .skip(offset)
            )
            
            # Convert ObjectIds and timestamps
            for file_record in files:
                file_record['_id'] = str(file_record['_id'])
                file_record['uploaded_at'] = file_record['uploaded_at'].isoformat()
                if 'processed_at' in file_record and file_record['processed_at']:
                    file_record['processed_at'] = file_record['processed_at'].isoformat()
                
                # Remove file path for security
                file_record.pop('file_path', None)
            
            total_count = db.uploaded_files.count_documents({'user_id': user_id})
            
            return {
                'success': True,
                'files': files,
                'total': total_count,
                'limit': limit,
                'offset': offset
            }
            
        except Exception as e:
            logger.error(f"Error getting user files: {str(e)}")
            return {'success': False, 'error': 'Failed to get files'}

    def delete_file(self, user_id: str, file_id: str) -> Dict:
        """Delete uploaded file"""
        try:
            db = db_instance.get_db()
            
            # Get file record
            file_record = db.uploaded_files.find_one({
                '_id': file_id,
                'user_id': user_id
            })
            
            if not file_record:
                return {'success': False, 'error': 'File not found'}
            
            # Delete physical file
            file_path = file_record.get('file_path')
            if file_path and os.path.exists(file_path):
                os.remove(file_path)
            
            # Delete database record
            db.uploaded_files.delete_one({'_id': file_id, 'user_id': user_id})
            
            return {'success': True, 'message': 'File deleted successfully'}
            
        except Exception as e:
            logger.error(f"File deletion error: {str(e)}")
            return {'success': False, 'error': 'Failed to delete file'}

    def cleanup_old_files(self, days_old: int = 30):
        """Cleanup files older than specified days (background task)"""
        try:
            db = db_instance.get_db()
            
            cutoff_date = datetime.utcnow() - pd.Timedelta(days=days_old)
            
            # Find old files
            old_files = db.uploaded_files.find({
                'uploaded_at': {'$lt': cutoff_date}
            })
            
            deleted_count = 0
            for file_record in old_files:
                try:
                    # Delete physical file
                    file_path = file_record.get('file_path')
                    if file_path and os.path.exists(file_path):
                        os.remove(file_path)
                    
                    # Delete database record
                    db.uploaded_files.delete_one({'_id': file_record['_id']})
                    deleted_count += 1
                    
                except Exception as e:
                    logger.error(f"Error deleting file {file_record.get('unique_filename')}: {str(e)}")
            
            logger.info(f"Cleaned up {deleted_count} old files")
            return {'success': True, 'deleted_count': deleted_count}
            
        except Exception as e:
            logger.error(f"File cleanup error: {str(e)}")
            return {'success': False, 'error': str(e)}content)
            
            return {
                'success': True,
                'content': content.strip(),
                'context': {
                    'file_type': 'pdf',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'word_count': len(content.split()),
                    'language': self._detect_language(content)
                }
            }
            
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            return {'success': False, 'error': f'PDF processing failed: {str(e)}'}

    def _process_document(self, file_path: str) -> Dict:
        """Process Word documents"""
        try:
            if file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                
                # Extract text from paragraphs
                content_parts = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        content_parts.append(paragraph.text.strip())
                
                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            content_parts.append(" | ".join(row_text))
                
                content = "\n\n".join(content_parts)
                
                # Metadata
                properties = doc.core_properties
                metadata = {
                    'title': properties.title or '',
                    'author': properties.author or '',
                    'subject': properties.subject or '',
                    'created': properties.created.isoformat() if properties.created else '',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables)
                }
                
            else:
                # Fallback for .doc files using textract
                content = textract.process(file_path).decode('utf-8')
                metadata = {'format': 'doc', 'processed_with': 'textract'}
            
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(content)
            
            return {
                'success': True,
                'content': content.strip(),
                'context': {
                    'file_type': 'document',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'word_count': len(content.split()),
                    'language': self._detect_language(content)
                }
            }
            
        except Exception as e:
            logger.error(f"Document processing error: {str(e)}")
            return {'success': False, 'error': f'Document processing failed: {str(e)}'}

    def _process_presentation(self, file_path: str) -> Dict:
        """Process PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            total_text = ""
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slide_content = "\n".join(slide_text)
                slides_content.append(f"Slide {i+1}:\n{slide_content}")
                total_text += slide_content + "\n\n"
            
            content = "\n\n---\n\n".join(slides_content)
            
            metadata = {
                'slides_count': len(prs.slides),
                'slide_layouts': len(prs.slide_layouts),
                'slide_masters': len(prs.slide_masters)
            }
            
            summary = self._generate_content_summary(total_text)
            key_points = self._extract_key_points(total_text)
            
            return {
                'success': True,
                'content': content.strip(),
                'context': {
                    'file_type': 'presentation',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'slides_count': len(prs.slides),
                    'language': self._detect_language(total_text)
                }
            }
            
        except Exception as e:
            logger.error(f"Presentation processing error: {str(e)}")
            return {'success': False, 'error': f'Presentation processing failed: {str(e)}'}

    def _process_image(self, file_path: str) -> Dict:
        """Process images using OCR"""
        try:
            image = Image.open(file_path)
            
            # Image metadata
            metadata = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'has_transparency': image.mode in ('RGBA', 'LA') or 'transparency' in image.info
            }
            
            # Extract text using OCR
            try:
                content = pytesseract.image_to_string(image)
                content = content.strip()
            except:
                content = ""
            
            if not content:
                return {
                    'success': True,
                    'content': "",
                    'context': {
                        'file_type': 'image',
                        'metadata': metadata,
                        'summary': "Image file with no extractable text",
                        'key_points': [],
                        'has_text': False
                    }
                }
            
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(content)
            
            return {
                'success': True,
                'content': content,
                'context': {
                    'file_type': 'image',
                    'metadata': metadata,
                    'summary': summary,
                    'key_points': key_points,
                    'has_text': True,
                    'language': self._detect_language(content)
                }
            }
            
        except Exception as e:
            logger.error(f"Image processing error: {str(e)}")
            return {'success': False, 'error': f'Image processing failed: {str(e)}'}

    def _process_text_file(self, file_path: str) -> Dict:
        """Process text files"""
        try:
            # Detect encoding
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if not content:
                return {'success': False, 'error': 'Could not decode text file'}
            
            metadata = {
                'encoding': encoding,
                'lines': len(content.split('\n')),
                'file_size': os.path.getsize(file_path)
            }
            
            summary = self._generate_content_summary(content)
            key_points = self._extract_key_points(